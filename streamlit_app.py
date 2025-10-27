import os                         # Standard library: environment variables, file checks, cleanup [web:61]
import tempfile                  # Create safe temporary files for uploaded content [web:61]
import base64                    # Reserved for potential binary/text encoding tasks (currently unused) [web:61]
import streamlit as st           # UI framework for the web app and chat elements [web:2]
from dotenv import load_dotenv   # Load environment variables from a .env file [web:61]

# LangChain document loaders for multiple file types [web:15]
from langchain_community.document_loaders import (
    PyPDFLoader,                  # PDF loader that returns page-wise Documents [web:15]
    TextLoader,                   # Plain text file loader [web:15]
    UnstructuredWordDocumentLoader,   # DOCX loader via unstructured [web:15]
    UnstructuredMarkdownLoader        # Markdown loader via unstructured [web:15]
)

from langchain.text_splitter import RecursiveCharacterTextSplitter  # Chunk text with overlaps for RAG [web:33]
from langchain_community.embeddings import HuggingFaceEmbeddings    # Sentence transformers embeddings [web:15]
from langchain_community.vectorstores import FAISS                  # FAISS vector index via LangChain [web:15]
from langchain_together import ChatTogether                         # Together AI chat model integration [web:21]
from pptx import Presentation                                       # Parse PPTX slides and shapes [web:32]
from fpdf import FPDF                                               # Lightweight PDF generation [web:7]
import io                                                           # In-memory byte and text buffers [web:61]

# ---------------- SETUP ----------------
load_dotenv()                                        # Load environment variables from .env into process env [web:61]
together_api_key = os.getenv("TOGETHER_API_KEY")     # Read Together API key for model access [web:23]
st.set_page_config(page_title="Tiet-Genie 🤖", layout="wide")  # Configure Streamlit page (title, wide layout) [web:41]

# Custom styling for better text visibility across chat and markdown [web:2]
st.markdown("""
<style>
.stChatMessageContent, .stMarkdown {
    color: #111 !important;
    font-weight: 500;
}
</style>
""", unsafe_allow_html=True)

# ---------------- SIDEBAR ----------------
with st.sidebar:                                                # Sidebar container for controls and exports [web:2]
    if os.path.exists("TIETlogo.png"):                          # Show logo if present in working directory [web:61]
        st.image("TIETlogo.png", width=120)                     # Display institute logo in sidebar [web:2]
    else:
        st.markdown("## 🏫 TIET")                               # Fallback title if logo is missing [web:2]
    
    st.markdown("## 🤖 Tiet-Genie")                             # App identity header [web:2]
    st.markdown("How can I assist you today? 😊")               # Friendly prompt in sidebar [web:2]
    uploaded_files = st.file_uploader(                          # Multi-file uploader for various doc types [web:47]
        "📎 Upload PDFs, DOCX, PPTX, TXT, or MD",
        type=["pdf", "docx", "pptx", "txt", "md"],
        accept_multiple_files=True
    )

# ---------------- LOAD DEFAULT PDFs ----------------
@st.cache_resource(show_spinner="Loading default PDFs...")       # Cache heavy resources (index) across reruns/sessions [web:49]
def load_default_vectorstore():
    default_files = ["rules.pdf", "AcademicRegulations.pdf"]     # Seed documents to preload if available [web:15]
    docs = []                                                    # Accumulate LangChain Document objects [web:15]
    
    for path in default_files:                                   # Iterate seed files defensively [web:61]
        if os.path.exists(path):                                 # Only attempt load if the file is present [web:61]
            try:
                docs.extend(PyPDFLoader(path).load())            # Load PDF pages into Documents [web:15]
            except Exception as e:
                st.warning(f"⚠️ Failed to load {path}: {str(e)}")# UI feedback for partial load failures [web:2]
        else:
            st.warning(f"⚠️ File not found: {path}")             # Inform missing default assets [web:2]
    
    if not docs:                                                 # If nothing loaded, initialize with a placeholder [web:15]
        st.info("ℹ️ No default PDFs found. Upload your own files to get started.")  # Guidance for first-time users [web:2]
        from langchain.docstore.document import Document         # Import only when needed to avoid overhead [web:15]
        dummy_doc = Document(
            page_content="Welcome to Tiet-Genie! Please upload your documents to get started.",
            metadata={}
        )
        docs = [dummy_doc]                                       # Basic content so FAISS can still be built [web:15]

    splitter = RecursiveCharacterTextSplitter(                   # Configure chunking granularity and overlap [web:33]
        chunk_size=500, chunk_overlap=50
    )
    chunks = splitter.split_documents(docs)                      # Page-to-chunk expansion for retrieval [web:33]
    embed = HuggingFaceEmbeddings(model_name="sentence-transformers/all-MiniLM-L6-v2")  # Compact, fast embeddings [web:15]
    return FAISS.from_documents(chunks, embed)                   # Build FAISS index from chunks + embeddings [web:15]

# Initialize vector store with error handling and early abort if it fails [web:2]
try:
    vector_store = load_default_vectorstore()                     # Cached FAISS index for default content [web:49]
except Exception as e:
    st.error(f"❌ Error initializing vector store: {str(e)}")      # Surface critical failures to the user [web:2]
    st.stop()                                                     # Halt app execution gracefully [web:2]

# ---------------- HANDLE USER FILE UPLOADS ----------------
def load_file_to_docs(file_path, ext):
    """
    Normalize various file types to a list of LangChain Document objects for a uniform ingest pipeline. [web:15]
    """
    try:
        if ext == "pdf":
            return PyPDFLoader(file_path).load()                  # PDF -> Documents (per page) [web:15]
        elif ext == "docx":
            return UnstructuredWordDocumentLoader(file_path).load()  # DOCX via unstructured [web:15]
        elif ext == "pptx":
            prs = Presentation(file_path)                         # Open the PowerPoint file [web:32]
            text = "\n".join(                                     # Extract visible text from shapes across slides [web:32]
                shape.text for slide in prs.slides
                for shape in slide.shapes if hasattr(shape, "text")
            )
            temp_txt = file_path + ".txt"                         # Persist as temp .txt for TextLoader [web:61]
            with open(temp_txt, "w", encoding="utf-8") as f:
                f.write(text)                                     # Write extracted PPTX text [web:32]
            return TextLoader(temp_txt).load()                    # Load text into Documents [web:15]
        elif ext == "txt":
            return TextLoader(file_path, encoding="utf-8").load() # Plain text loader with encoding [web:15]
        elif ext == "md":
            return UnstructuredMarkdownLoader(file_path).load()   # Markdown via unstructured [web:15]
        return []                                                 # Unknown extension -> no docs [web:61]
    except Exception as e:
        st.error(f"❌ Error loading file {file_path}: {str(e)}")   # Report ingestion errors per-file [web:2]
        return []                                                 # Fail soft to keep app responsive [web:61]

if uploaded_files:                                                # If user uploaded anything, process now [web:47]
    new_docs = []                                                 # Collected Documents from all uploads [web:15]
    for f in uploaded_files:                                      # Each uploaded file-like object from Streamlit [web:47]
        ext = f.name.split(".")[-1].lower()                       # Infer by file extension [web:61]
        try:
            with tempfile.NamedTemporaryFile(delete=False, suffix=f".{ext}") as tmp:
                tmp.write(f.read())                               # Write upload bytes to a temp file path [web:61]
                tmp_path = tmp.name                               # Capture path for loaders that require filenames [web:61]
            
            loaded_docs = load_file_to_docs(tmp_path, ext)        # Normalize into Documents [web:15]
            if loaded_docs:
                new_docs.extend(loaded_docs)                      # Accumulate across files [web:15]
                st.success(f"✅ Successfully loaded {f.name}")     # Positive feedback per file [web:2]
            else:
                st.warning(f"⚠️ No content extracted from {f.name}") # Inform empty extraction [web:2]
                
            try:
                os.unlink(tmp_path)                               # Clean up temp file if possible [web:61]
            except:
                pass                                              # If cleanup fails, continue safely [web:61]
                
        except Exception as e:
            st.error(f"❌ Error processing {f.name}: {str(e)}")    # Upload-level failure reporting [web:2]

    if new_docs:                                                  # Only rebuild when there is new content [web:15]
        try:
            splitter = RecursiveCharacterTextSplitter(            # Reuse same chunking strategy for consistency [web:33]
                chunk_size=500, chunk_overlap=50
            )
            new_chunks = splitter.split_documents(new_docs)       # Chunk user-provided documents [web:33]
            embed = HuggingFaceEmbeddings(model_name="sentence-transformers/all-MiniLM-L6-v2")  # Same embedding model [web:15]
            new_vs = FAISS.from_documents(new_chunks, embed)      # Build a temporary FAISS index [web:15]
            vector_store.merge_from(new_vs)                       # Merge into main index to expand knowledge base [web:15]
            st.success(f"✅ Added {len(new_docs)} documents to knowledge base")  # Confirm ingest count [web:2]
        except Exception as e:
            st.error(f"❌ Error updating vector store: {str(e)}") # Index rebuild failures surfaced to UI [web:2]

# ---------------- LLM + RETRIEVER ----------------
try:
    retriever = vector_store.as_retriever(                        # Wrap FAISS as a retriever interface [web:15]
        search_type="mmr",                                        # Maximal Marginal Relevance for diversity [web:15]
        search_kwargs={"k": 4, "fetch_k": 10, "lambda_mult": 0.5} # Retrieve diverse top-k from a larger pool [web:15]
    )

    if not together_api_key:                                      # Validate presence of API key before LLM calls [web:23]
        st.error("❌ TOGETHER_API_KEY not found in environment variables. Please set it in your .env file.")
        st.stop()                                                 # Stop to prevent unauthenticated requests [web:2]

    llm = ChatTogether(                                           # Instantiate Together chat client via LangChain [web:21]
        model="deepseek-ai/DeepSeek-V3",                          # Model served by Together AI [web:21]
        temperature=0.2,                                          # Lower temperature for focused, deterministic outputs [web:21]
        together_api_key=together_api_key                         # Provide credentials for calls [web:23]
    )
except Exception as e:
    st.error(f"❌ Error initializing LLM or retriever: {str(e)}")  # Report model/retriever failures [web:2]
    st.stop()                                                     # Abort since chat cannot function without LLM [web:2]

# ---------------- CHAT HISTORY ----------------
if "chat_history" not in st.session_state:                        # Persist conversation across reruns [web:2]
    st.session_state.chat_history = []                            # Initialize chat buffer [web:2]
if "greeted" not in st.session_state:                             # Track whether welcome banner was shown [web:2]
    st.session_state.greeted = False

if not st.session_state.greeted and not st.session_state.chat_history:
    # First-time greeting for an empty chat to improve UX [web:2]
    st.markdown("<h2 style='text-align:center;'>👋 Hello TIETian! How can I help you today?</h2>", unsafe_allow_html=True)

# ---------------- CHAT UI ----------------
for msg in st.session_state.chat_history:                         # Replay prior messages each rerun [web:1]
    with st.chat_message(msg["role"]):                            # Role-aware chat bubble (user/assistant) [web:1]
        st.markdown(msg["message"], unsafe_allow_html=True)       # Render rich text message content [web:2]

user_prompt = st.chat_input("Ask something about TIET or the lecture notes...")  # Bottom-aligned chat input [web:2]
if user_prompt:
    with st.chat_message("user"):                                 # Immediately echo user message in UI [web:1]
        st.write(user_prompt)                                     # Show plain user text [web:2]
    st.session_state.chat_history.append({"role": "user", "message": user_prompt})  # Store turn in history [web:2]

    with st.chat_message("assistant"):                            # Assistant response bubble [web:1]
        with st.spinner("Thinking..."):                           # Visual progress indicator during retrieval/LLM [web:2]
            try:
                retrieved_docs = retriever.get_relevant_documents(user_prompt)  # Search index for relevant chunks [web:15]

                # Build a compact context block from retrieved Documents (include page hints for source section) [web:15]
                context_text = "\n\n".join([
                    f"[Page {doc.metadata.get('page', '?')}] {doc.page_content.strip()}" for doc in retrieved_docs
                ])

                # System-style prompt instructing model to ground answers in provided snippets and avoid inline citations [web:21]
                prompt_to_llm = f"""
You are an AI assistant for Thapar Institute. Use the following document snippets to answer the question. Be specific and provide detailed information, but DO NOT include page numbers or citations like [Page X] in your main answer. The source snippets will be provided separately.

--- DOCUMENTS ---
{context_text}

--- QUESTION ---
{user_prompt}
"""

                response_obj = llm.invoke(prompt_to_llm)          # Call Together AI chat completion via LangChain [web:21]
                response = response_obj.content.strip() if hasattr(response_obj, "content") else str(response_obj).strip()  # Normalize text [web:21]

                # Heuristics to detect "no relevant info found" phrasings and suppress source section in that case [web:2]
                no_info_indicators = [
                    "do not contain specific information",
                    "do not contain information",
                    "does not contain specific information",
                    "does not contain information",
                    "do not address",
                    "does not address",
                    "do not discuss",
                    "does not discuss",
                    "no information about",
                    "no specific information about",
                    "not mentioned in the documents",
                    "not found in the documents",
                    "documents do not mention",
                    "excerpts do not contain",
                    "snippets do not contain",
                    "provided documents do not"
                ]
                
                response_lower = response.lower()                  # Case-normalize for matching [web:61]
                show_sources = not any(indicator in response_lower for indicator in no_info_indicators)  # Decide visibility [web:2]
                
                if show_sources:
                    # Render short source snippets only when relevant content was found to support the answer [web:2]
                    source_section = "\n\n---\n\n**📄 Source Snippets:**\n"
                    for i, doc in enumerate(retrieved_docs, 1):
                        page = doc.metadata.get("page", "?")       # Page metadata if available [web:15]
                        snippet = doc.page_content.strip().replace("\n", " ")[:300]  # Trim for readability [web:2]
                        source_section += f"- **Snippet {i} (Page {page})**: {snippet}\n"
                    final_response = f"{response}\n{source_section}"  # Append sources below the assistant text [web:2]
                else:
                    final_response = response                      # Keep output clean if nothing relevant was found [web:2]

                st.markdown(final_response, unsafe_allow_html=True)  # Display assistant response (and sources if any) [web:2]
                st.session_state.chat_history.append({"role": "assistant", "message": final_response})  # Persist assistant turn [web:2]

            except Exception as e:
                error_response = f"⚠️ Error: {str(e)}"            # Defensive error capture for retrieval/LLM calls [web:2]
                st.markdown(error_response)                        # Show error in the chat bubble [web:2]
                st.session_state.chat_history.append({"role": "assistant", "message": error_response})  # Persist error text [web:2]

# ---------------- EXPORT CHAT HISTORY ----------------
def export_chat_history():
    """
    Build downloadable PDF and TXT exports of the chat history using in-memory buffers and Streamlit buttons. [web:7]
    """
    chat = st.session_state.chat_history                          # Snapshot current conversation [web:2]
    if not chat:
        return                                                     # Nothing to export [web:61]

    try:
        # --- PDF Export (Using built-in fonts only) ---
        pdf = FPDF()                                               # Create new PDF document [web:7]
        pdf.add_page()                                             # Start first page [web:7]
        pdf.set_font("Arial", size=12)                             # Built-in core font for compatibility [web:7]
        pdf.set_auto_page_break(auto=True, margin=15)              # Automatic paging with margins [web:7]
        
        for msg in chat:
            role = "You" if msg["role"] == "user" else "Tiet-Genie"  # Human-readable role labels [web:2]
            clean_message = msg['message'].replace('**', '').replace('*', '').replace('#', '')  # Strip markdown [web:2]
            clean_message = clean_message.encode('latin1', 'ignore').decode('latin1')           # Avoid non-Latin1 glyph issues [web:7]
            pdf.multi_cell(0, 10, f"{role}:\n{clean_message}\n")  # Write wrapped text to PDF [web:7]

        # Output PDF into memory instead of disk for Streamlit download [web:7]
        pdf_output = pdf.output(dest='S')                          # 'S' -> return as string/bytes [web:7]
        if isinstance(pdf_output, str):
            pdf_bytes = pdf_output.encode('latin1')                # Normalize to bytes if string [web:7]
        else:
            pdf_bytes = pdf_output
        
        pdf_buffer = io.BytesIO(pdf_bytes)                         # Bytes buffer for download button [web:61]

        # --- TXT Export ---
        txt_buffer = io.StringIO()                                 # Text buffer for plain export [web:61]
        for msg in chat:
            role = "You" if msg["role"] == "user" else "Tiet-Genie"
            clean_txt = msg['message'].replace('**', '').replace('*', '').replace('#', '')  # Strip markdown [web:2]
            txt_buffer.write(f"{role}:\n{clean_txt}\n\n")          # Simple readable transcript [web:2]
        txt_buffer.seek(0)                                         # Rewind for reading [web:61]

        st.sidebar.markdown("### 📤 Export Chat History")          # Sidebar export section [web:2]
        st.sidebar.download_button(                                # Download PDF via in-memory buffer [web:47]
            "⬇️ Download as .pdf",
            data=pdf_buffer,
            file_name="chat_history.pdf",
            mime="application/pdf"
        )
        st.sidebar.download_button(                                # Download TXT transcript [web:47]
            "⬇️ Download as .txt",
            data=txt_buffer.getvalue(),
            file_name="chat_history.txt",
            mime="text/plain"
        )
        
    except Exception as e:
        st.sidebar.error(f"Export error: {str(e)}")                # Report export failures [web:2]
        # Fallback: Only TXT export if PDF fails
        try:
            txt_buffer = io.StringIO()                             # Ensure user still gets a transcript [web:61]
            for msg in chat:
                role = "You" if msg["role"] == "user" else "Tiet-Genie"
                clean_txt = msg['message'].replace('**', '').replace('*', '').replace('#', '')  # Strip markdown [web:2]
                txt_buffer.write(f"{role}:\n{clean_txt}\n\n")
            txt_buffer.seek(0)
            
            st.sidebar.markdown("### 📤 Export Chat History")      # Present fallback export option [web:2]
            st.sidebar.download_button(                            # TXT-only fallback [web:47]
                "⬇️ Download as .txt",
                data=txt_buffer.getvalue(),
                file_name="chat_history.txt",
                mime="text/plain"
            )
        except Exception as fallback_error:
            st.sidebar.error(f"Export failed: {str(fallback_error)}")  # Last-resort error path [web:2]

export_chat_history()                                              # Always render export controls if chat exists [web:2]
